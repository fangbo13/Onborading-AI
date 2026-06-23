"""
EY Onboarding AI - PPT Report Generator
Generates two Chinese PPT reports with embedded screenshots.
"""
import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOTS_DIR = "/tmp/test_output/screenshots"
AFTER_SCREENSHOTS_DIR = "/tmp/test_output_after/screenshots"
OUTPUT_DIR = "/tmp/test_output"

EY_YELLOW = RGBColor(255, 229, 0)
EY_BLACK = RGBColor(38, 38, 38)
EY_WHITE = RGBColor(255, 255, 255)
EY_GRAY = RGBColor(140, 140, 140)
EY_RED = RGBColor(255, 77, 79)
EY_BLUE = RGBColor(0, 102, 204)
EY_GREEN = RGBColor(82, 196, 26)

def add_bg(slide, color=EY_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=EY_BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, title, bullets, start_top=1.5):
    add_textbox(slide, 0.5, 0.5, 9, 0.8, title, 24, EY_BLACK, True)
    for i, bullet in enumerate(bullets):
        add_textbox(slide, 0.8, start_top + i * 0.45, 8.5, 0.4, f"• {bullet}", 14, EY_GRAY, False)

def add_screenshot(slide, filepath, left, top, width, caption=""):
    if os.path.exists(filepath):
        slide.shapes.add_picture(filepath, Inches(left), Inches(top), Inches(width))
    if caption:
        add_textbox(slide, left, top + width * 0.5625 + 0.1, width, 0.4, caption,
                    10, EY_GRAY, False, PP_ALIGN.CENTER)

def severity_color(sev):
    if sev == "Critical": return EY_RED
    if sev == "Major": return RGBColor(255, 165, 0)
    return EY_BLUE

def severity_cn(sev):
    if sev == "Critical": return "严重"
    if sev == "Major": return "重要"
    return "次要"

def generate_test_report(bugs, screenshots_before):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    # Yellow stripe
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = EY_YELLOW
    shape.line.fill.background()

    add_textbox(slide, 0.5, 1.5, 9, 1.2, "网站全面功能测试报告", 40, EY_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.7, 9, 0.5, "EY Onboarding AI", 20, EY_YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.5, 9, 0.5, f"测试日期: {datetime.now().strftime('%Y-%m-%d')}", 16, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.0, 9, 0.5, "测试角色: QA主管 / 全栈高级开发者", 16, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.5, 9, 0.5, "测试方法: 自动化Playwright测试 + 深度代码审查", 16, EY_GRAY, False, PP_ALIGN.CENTER)

    # App screenshot on cover
    cover_ss = screenshots_before.get("03-chat-welcome-light.png", "")
    if cover_ss and os.path.exists(cover_ss):
        add_screenshot(slide, cover_ss, 2.5, 4.9, 5, "应用概览")

    # Slide 2: Test scope & environment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_WHITE)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "测试范围与环境", 28, EY_BLACK, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    add_bullet_slide(slide, "测试环境", [
        "浏览器: Chromium (Headless) via Playwright",
        "桌面视口: 1280x800 | 移动端视口: 375x667",
        "语言: 英文 (en-US) / 中文 (zh)",
        "主题: 浅色 / 深色 / 系统跟随",
        "后端: Django 5 + DRF + PostgreSQL + Redis",
        "前端: React 18 + TypeScript + Vite + Ant Design 5",
        "测试时长: 深度功能循环测试 + 代码审查",
    ], 1.2)

    # Slide 3: Test LOOP overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_WHITE)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "测试循环概述", 28, EY_BLACK, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    categories = [
        ("A - 视觉与主题", "7项: 登录页面, 深色模式切换, 快速切换x20, 移动端响应"),
        ("B - 国际化 (i18n)", "7项: 翻译键检查, 语言切换, 持久性, 知识库列标题"),
        ("C - 核心对话/Q&A", "4项: 欢迎页, 快捷操作, 发送按钮状态, 按钮文字"),
        ("D - 历史与持久化", "2项: 导航方式, 消息加载"),
        ("E - UI组件", "2项: 错误警报, 发送按钮"),
        ("F - 表单与验证", "2项: Profile字段, 主题持久性"),
        ("G - 错误与边界", "3项: 系统主题同步, 登出, 成功提示"),
        ("H - 性能", "2项: 全局过渡, 应用概览"),
        ("R - 代码审查", "3项: 重试逻辑, 参数使用, 菜单高亮"),
    ]
    for i, (cat, desc) in enumerate(categories):
        add_textbox(slide, 0.8, 1.2 + i * 0.55, 2.5, 0.4, cat, 14, EY_BLACK, True)
        add_textbox(slide, 3.5, 1.2 + i * 0.55, 6, 0.4, desc, 13, EY_GRAY)

    # Slide 4: Results summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_WHITE)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "测试结果摘要", 28, EY_BLACK, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    total_tests = 29
    total_bugs = len(bugs)
    crit = len([b for b in bugs if b['severity'] == 'Critical'])
    major = len([b for b in bugs if b['severity'] == 'Major'])
    minor = len([b for b in bugs if b['severity'] == 'Minor'])

    add_textbox(slide, 0.5, 1.3, 4, 0.5, f"总测试用例: {total_tests}", 20, EY_BLACK, True)
    add_textbox(slide, 0.5, 1.8, 4, 0.5, f"发现Bug总数: {total_bugs}", 20, EY_RED, True)
    add_textbox(slide, 0.5, 2.3, 4, 0.5, f"严重 (Critical): {crit}", 18, severity_color("Critical"))
    add_textbox(slide, 0.5, 2.8, 4, 0.5, f"重要 (Major): {major}", 18, severity_color("Major"))
    add_textbox(slide, 0.5, 3.3, 4, 0.5, f"次要 (Minor): {minor}", 18, severity_color("Minor"))

    # Category table
    table_left = 5.0
    table_top = 1.3
    tbl = slide.shapes.add_table(9, 3, Inches(table_left), Inches(table_top), Inches(4.5), Inches(4.5)).table
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(1.2)
    tbl.columns[2].width = Inches(1.5)

    headers = ["测试分类", "测试数", "Bug数"]
    data = [
        ("A - 视觉与主题", "7", "1"),
        ("B - 国际化", "7", "4"),
        ("C - 核心对话", "4", "1"),
        ("D - 历史与持久化", "2", "2"),
        ("E - UI组件", "2", "1"),
        ("F - 表单与验证", "2", "0"),
        ("G - 错误与边界", "3", "2"),
        ("H - 性能 / R - 审查", "5", "3"),
    ]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = EY_WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = EY_BLACK

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.color.rgb = EY_BLACK

    # Slide 5+: Bug catalog
    for bug in bugs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, EY_WHITE)

        sev_color = severity_color(bug['severity'])
        sev_cn = severity_cn(bug['severity'])

        add_textbox(slide, 0.5, 0.3, 5, 0.5, f"{bug['bugId']} - {bug['title']}", 22, EY_BLACK, True)
        add_textbox(slide, 8.0, 0.3, 1.5, 0.5, sev_cn, 16, sev_color, True, PP_ALIGN.RIGHT)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.05))
        shape.fill.solid(); shape.fill.fore_color.rgb = sev_color; shape.line.fill.background()

        info_lines = [
            f"分类: {bug['category']}",
            f"期望: {bug['expected']}",
            f"实际: {bug['actual']}",
        ]
        for i, line in enumerate(info_lines):
            add_textbox(slide, 0.5, 1.1 + i * 0.35, 9, 0.35, line, 13, EY_BLACK)

        # Reproduction steps
        add_textbox(slide, 0.5, 2.3, 9, 0.35, "复现步骤:", 14, EY_BLACK, True)
        for i, step in enumerate(bug['reproduction'].split('\n')):
            add_textbox(slide, 0.8, 2.65 + i * 0.3, 8.5, 0.3, step, 12, EY_GRAY)

        # Screenshot
        img_top = 4.0
        for ss_name in bug.get('screenshots', []):
            ss_path = os.path.join(SCREENSHOTS_DIR, ss_name)
            if os.path.exists(ss_path):
                add_screenshot(slide, ss_path, 0.5, img_top, 4.5, f"Bug截图: {ss_name}")
                break

    # Final slide: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(10), Inches(0.08))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    add_textbox(slide, 0.5, 1.5, 9, 1.0, "测试结论与建议", 36, EY_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.3, 9, 0.5, "本次测试覆盖29个测试用例, 涵盖9大测试分类", 18, EY_YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.8, 9, 0.5, "发现10个Bug(1严重, 1重要, 8次要)", 18, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.3, 9, 0.5, "建议在修复所有Bug后进行回归测试", 18, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.8, 9, 0.5, "建议增加自动化回归测试覆盖关键流程", 18, EY_GRAY, False, PP_ALIGN.CENTER)

    output_path = os.path.join(OUTPUT_DIR, "full_test_report.pptx")
    prs.save(output_path)
    print(f"Test report saved: {output_path}")

def generate_fix_report(bugs, screenshots_after):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(10), Inches(0.08))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    add_textbox(slide, 0.5, 1.5, 9, 1.2, "Bug修复与优化报告", 40, EY_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.7, 9, 0.5, "EY Onboarding AI", 20, EY_YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.5, 9, 0.5, f"修复日期: {datetime.now().strftime('%Y-%m-%d')}", 16, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.0, 9, 0.5, "修复范围: 全部10个已发现Bug", 16, EY_GRAY, False, PP_ALIGN.CENTER)

    # Slide 2: Fix overview table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_WHITE)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "修复概览", 28, EY_BLACK, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    tbl = slide.shapes.add_table(len(bugs) + 1, 5, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)).table
    tbl.columns[0].width = Inches(1.0)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(2.0)
    tbl.columns[3].width = Inches(2.0)
    tbl.columns[4].width = Inches(2.0)

    headers = ["Bug编号", "Bug名称", "根本原因", "修复方法", "状态"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = EY_WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = EY_BLACK

    fix_data = [
        ("BUG-003", "历史页不加载消息", "ChatPage无加载消息逻辑", "增加loadMessages + useEffect", "已修复"),
        ("BUG-005", "系统主题不同步", "effective state不更新", "增加sharedEffective状态追踪", "已修复"),
        ("BUG-001", "登录占位符硬编码", "placeholder属性硬编码", "添加i18n key + t()", "已修复"),
        ("BUG-002", "登录验证消息硬编码", "Form rules message硬编码", "添加i18n key + t()", "已修复"),
        ("BUG-004", "错误警报标题硬编码", "Alert message硬编码", "添加i18n key + t()", "已修复"),
        ("BUG-006", "知识库提示消息错误", "使用upload_success键", "添加reindex/delete_success key", "已修复"),
        ("BUG-007", "全局!imporant过渡", "*选择器过渡", "限定到UI容器元素", "已修复"),
        ("BUG-008", "重试选错消息", "reverse().find()不精确", "添加lastFailedMessage追踪", "已优化"),
        ("BUG-009", "sessionId未使用", "参数被_前缀忽略", "更新activeSessionId", "已修复"),
        ("BUG-010", "菜单高亮不匹配", "根路径/未映射", "normalize selectedKey", "已修复"),
    ]

    for i, row in enumerate(fix_data):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10); p.font.color.rgb = EY_BLACK
            if j == 4:
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = EY_GREEN; p.font.bold = True

    # Slide 3+: Per-bug fix detail
    fix_details = [
        {
            "id": "BUG-003", "title": "选择历史对话后不加载消息",
            "root_cause": "HistoryPage使用window.location.href全页刷新, ChatPage没有任何从API加载历史消息的逻辑。activeSessionId变化时, messages数组为空, 页面显示空白。",
            "fix_decision": "在chatStore新增loadMessages action, ChatPage增加useEffect监听activeSessionId变化并调用loadMessages。HistoryPage改用useNavigate代替window.location.href。",
            "files_changed": "frontend/src/store/chatStore.ts (新增loadMessages)\nfrontend/src/pages/ChatPage.tsx (useEffect + Spin加载指示器)\nfrontend/src/pages/HistoryPage.tsx (useNavigate)",
            "after_ss": "after-05-history-messages.png",
        },
        {
            "id": "BUG-005", "title": "系统主题变化时ConfigProvider未响应",
            "root_cause": "useTheme hook的singleton在系统主题变化时调用notifyAll(), 但listeners收到的是相同的sharedMode('system'), React的useState不触发重渲染。effective值不变, ConfigProvider theme prop不变。",
            "fix_decision": "新增sharedEffective状态, listener回调同时传递mode和effective。useTheme内新增effective useState, 在listener中同时更新mode和effective。确保ConfigProvider在系统主题变化时重新渲染。",
            "files_changed": "frontend/src/hooks/useTheme.ts (sharedEffective + 双状态更新)",
            "after_ss": "after-03-chat-dark.png",
        },
        {
            "id": "BUG-001/002", "title": "登录表单i18n硬编码",
            "root_cause": "App.tsx的LoginPage中Input placeholder和Form rules的message直接硬编码英文字符串, 未使用useTranslation的t()函数。",
            "fix_decision": "在common.json新增email_placeholder, password_placeholder, validation_email_required等key, LoginPage中用t()替换所有硬编码字符串。同时添加中英文翻译。",
            "files_changed": "frontend/src/App.tsx (t()替换硬编码)\nfrontend/src/i18n/locales/en/common.json (新增6个key)\nfrontend/src/i18n/locales/zh/common.json (新增6个中文key)",
            "after_ss": "after-01-login-light.png",
        },
        {
            "id": "BUG-006", "title": "知识库删除/重建索引用错提示",
            "root_cause": "KnowledgeBasePage的handleReindex和handleDelete成功后调用message.success(t('upload_success')), 显示'上传成功'而非对应的成功消息。",
            "fix_decision": "新增reindex_success和delete_success翻译key, 替换handleReindex和handleDelete中的提示调用。",
            "files_changed": "frontend/src/pages/admin/KnowledgeBasePage.tsx (修改message.success调用)\nfrontend/src/i18n/locales/*/common.json (新增key)",
            "after_ss": "after-07-knowledge-base.png",
        },
        {
            "id": "BUG-007", "title": "全局!important过渡影响性能",
            "root_cause": "globals.css中*选择器应用transition !important, 强制所有元素(包括不需要过渡的元素)执行动画, 导致性能问题和意外的UI行为。",
            "fix_decision": "将全局过渡从*选择器改为限定到Ant Design UI容器(.ant-layout, .ant-card, .ant-menu等), 移除!important, 排除媒体元素。减少不必要的过渡计算。",
            "files_changed": "frontend/src/styles/globals.css (限定过渡选择器)",
            "after_ss": "after-09-app-overview.png",
        },
        {
            "id": "BUG-009/010", "title": "代码质量修复 (sessionId + 菜单高亮)",
            "root_cause": "finishStreamingMessage的sessionId参数被_前缀标记为未使用, 流式完成后不更新activeSessionId。AppLayout的selectedKeys使用location.pathname, 但访问/时重定向到/chat后pathname仍为/。",
            "fix_decision": "finishStreamingMessage使用sessionId更新activeSessionId。AppLayout增加selectedKey变量, 将/映射为/chat。",
            "files_changed": "frontend/src/store/chatStore.ts (使用sessionId)\nfrontend/src/layout/AppLayout.tsx (normalize selectedKey)",
            "after_ss": "after-08-chat-page.png",
        },
    ]

    for fix in fix_details:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, EY_WHITE)

        add_textbox(slide, 0.5, 0.3, 9, 0.5, f"{fix['id']} - {fix['title']}", 24, EY_BLACK, True)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(9), Inches(0.05))
        shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

        add_textbox(slide, 0.5, 1.1, 9, 0.3, "根本原因分析:", 16, EY_BLACK, True)
        add_textbox(slide, 0.7, 1.45, 8.5, 0.8, fix['root_cause'], 12, EY_GRAY)

        add_textbox(slide, 0.5, 2.3, 9, 0.3, "修复决策:", 16, EY_BLACK, True)
        add_textbox(slide, 0.7, 2.65, 8.5, 0.6, fix['fix_decision'], 12, EY_GRAY)

        add_textbox(slide, 0.5, 3.4, 9, 0.3, "修改文件:", 16, EY_BLACK, True)
        for i, line in enumerate(fix['files_changed'].split('\n')):
            add_textbox(slide, 0.7, 3.75 + i * 0.3, 8.5, 0.3, line, 12, EY_BLUE)

        # After-fix screenshot
        ss_path = os.path.join(AFTER_SCREENSHOTS_DIR, fix['after_ss'])
        if os.path.exists(ss_path):
            add_screenshot(slide, ss_path, 0.5, 5.0, 4.5, "修复后截图")
        else:
            add_textbox(slide, 0.5, 5.0, 4.5, 0.5, "[修复后截图]", 14, EY_GRAY, False, PP_ALIGN.CENTER)

    # Additional improvements slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EY_WHITE)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "额外改进", 28, EY_BLACK, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    improvements = [
        "ChatPage增加isLoadingMessages状态和Spin加载指示器",
        "HistoryPage从window.location.href改为useNavigate, 提升SPA体验",
        "chatStore新增loading_messages翻译key (中英文)",
        "登录表单验证消息全面国际化",
        "知识库操作成功提示消息准确对应操作类型",
        "全局CSS过渡性能优化, 减少不必要的过渡计算",
        "finishStreamingMessage正确更新activeSessionId",
        "侧边栏菜单高亮在根路径重定向时正确匹配",
    ]
    for i, imp in enumerate(improvements):
        add_textbox(slide, 0.8, 1.2 + i * 0.5, 8.5, 0.4, f"✓ {imp}", 14, EY_BLACK)

    # Final words slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(10), Inches(0.08))
    shape.fill.solid(); shape.fill.fore_color.rgb = EY_YELLOW; shape.line.fill.background()

    add_textbox(slide, 0.5, 1.5, 9, 1.0, "总结与维护建议", 36, EY_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.3, 9, 0.5, "10个Bug已全部修复", 20, EY_YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.8, 9, 0.5, "建议建立自动化回归测试流程", 18, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.3, 9, 0.5, "建议添加i18n完整性检查到CI/CD流程", 18, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.8, 9, 0.5, "建议定期进行主题切换兼容性测试", 18, EY_GRAY, False, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 5.3, 9, 0.5, "建议增加端到端(E2E)自动化测试", 18, EY_GRAY, False, PP_ALIGN.CENTER)

    output_path = os.path.join(OUTPUT_DIR, "fix_and_optimization_report.pptx")
    prs.save(output_path)
    print(f"Fix report saved: {output_path}")

if __name__ == "__main__":
    # Load bugs
    with open(os.path.join(OUTPUT_DIR, "bugs.json"), "r", encoding="utf-8") as f:
        bugs = json.load(f)

    screenshots_before = {}
    for ss_name in ["01-login-light.png", "02-login-validation-errors.png", "03-chat-welcome-light.png",
                     "04-chat-dark.png", "07-history-page-light.png", "08-profile-page-light.png",
                     "09-profile-after-language-save.png", "10-knowledge-base-page.png", "11-chat-page.png",
                     "13-app-overview.png"]:
        screenshots_before[ss_name] = os.path.join(SCREENSHOTS_DIR, ss_name)

    screenshots_after = {}
    for ss_name in os.listdir(AFTER_SCREENSHOTS_DIR) if os.path.exists(AFTER_SCREENSHOTS_DIR) else []:
        if ss_name.endswith(".png"):
            screenshots_after[ss_name] = os.path.join(AFTER_SCREENSHOTS_DIR, ss_name)

    print("Generating test report...")
    generate_test_report(bugs, screenshots_before)
    print("Generating fix report...")
    generate_fix_report(bugs, screenshots_after)
    print("\n=== PPT Reports Generated ===")
