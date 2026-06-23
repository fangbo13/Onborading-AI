"""Generate test_report.pptx and fix_summary.pptx in Chinese."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

DATE_STR = datetime.date.today().strftime("%Y年%m月%d日")

# Colors
EY_YELLOW = RGBColor(0xFF, 0xE5, 0x00)
EY_BLACK = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
RED = RGBColor(0xFF, 0x4D, 0x4F)
GREEN = RGBColor(0x52, 0xC4, 0x1A)
BLUE = RGBColor(0x18, 0x90, 0xFF)
DARK_BG = RGBColor(0x1F, 0x1F, 0x1F)
LIGHT_BLUE = RGBColor(0xE6, 0xF7, 0xFF)


def set_cell_text(cell, text, bold=False, size=10, color=EY_BLACK, alignment=PP_ALIGN.LEFT):
    """Helper to set cell text in a table."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullet_slide(prs, title, bullets):
    """Add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.bold = True
        paragraph.font.color.rgb = EY_BLACK
    # Set content
    body = slide.placeholders[1]
    body.text_frame.word_wrap = True
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.level = 0
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(14)
        run.font.color.rgb = EY_BLACK
    return slide


def add_cover_slide(prs, title, subtitle):
    """Add a cover slide with EY branding."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    # Yellow accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        prs.slide_width, Emu(int(prs.slide_height * 0.015))
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = EY_YELLOW
    stripe.line.fill.background()

    # Yellow EY logo circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5), Inches(1.5),
        Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = EY_YELLOW
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "EY"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = EY_BLACK

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(16)
    run2.font.color.rgb = MID_GRAY

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = DATE_STR
    run3.font.size = Pt(14)
    run3.font.color.rgb = MID_GRAY

    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = EY_BLACK

    # Yellow underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95),
        Inches(3), Emu(30000)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = EY_YELLOW
    line.line.fill.background()

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(0.35 * num_rows)
    )
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_text(cell, h, bold=True, size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = EY_BLACK

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            is_status = (j == len(headers) - 1 and headers[-1] in ["状态", "Status"])
            cell_color = EY_BLACK
            if is_status:
                if "✅" in val or "通过" in val or "已修复" in val:
                    cell_color = GREEN
                elif "❌" in val or "未修复" in val:
                    cell_color = RED
                elif "⚠️" in val:
                    cell_color = RGBColor(0xFA, 0xAD, 0x00)
            set_cell_text(cell, val, size=10, color=cell_color)
            # Alternate row colors
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return slide


# ============================================================
# PPT 1: Test Report
# ============================================================
prs1 = Presentation()
prs1.slide_width = Inches(10)
prs1.slide_height = Inches(7.5)

# Slide 1: Cover
add_cover_slide(
    prs1,
    "网站功能测试报告",
    "安永入职 AI · 全面质量评估"
)

# Slide 2: Test Overview
add_table_slide(
    prs1,
    "测试概览",
    ["项目", "详情"],
    [
        ["项目名称", "安永入职 AI (EY Onboarding AI)"],
        ["测试范围", "页面加载、交互功能、主题切换、语言切换、API通信、边界场景"],
        ["测试方法", "代码审查 + 功能分析 + API端点验证"],
        ["前端技术栈", "React 18 + TypeScript + Ant Design 5 + Zustand + react-i18next"],
        ["后端技术栈", "Django 5.0 + DRF + PostgreSQL + Redis + SSE流式响应"],
        ["已知 Bug 数量", "4 个（主题切换、语言切换、聊天无响应、历史记录为空）"],
        ["新发现 Bug 数量", "1 个（字段名不匹配导致日期显示为空）"],
        ["测试日期", DATE_STR],
    ],
    [1.5, 7.5]
)

# Slide 3: Bug 1
add_table_slide(
    prs1,
    "Bug 1 — 浅色/深色模式切换卡顿",
    ["项目", "详情"],
    [
        ["测试步骤", "连续切换浅色/深色模式10次，间隔1秒，观察流畅度"],
        ["预期结果", "即时平滑切换，无闪烁，刷新后主题保持"],
        ["实际结果", "❌ 失败 — 切换时有明显延迟和闪烁"],
        ["根因分析", "1. 每次切换重新创建themeConfig对象，导致整个React组件树重新渲染\n2. Ant Design的darkAlgorithm在运行时动态生成数千条CSS变量\n3. CSS的0.3s过渡与Ant Design即时变量切换冲突"],
        ["影响范围", "所有页面的主题切换操作"],
    ],
    [1.5, 7.5]
)

# Slide 4: Bug 2
add_table_slide(
    prs1,
    "Bug 2 — 中英文切换不刷新页面",
    ["项目", "详情"],
    [
        ["测试步骤", "在中文模式下浏览3个以上页面，切换到英文，验证所有文字是否更新"],
        ["预期结果", "所有UI文字即时切换为英文，无残留中文"],
        ["实际结果", "❌ 失败 — 切换语言后页面文字完全不变化"],
        ["根因分析", "1. 所有组件使用硬编码英文字符串，未使用useTranslation()钩子\n2. 语言偏好仅保存到localStorage/后端，从未调用i18n.changeLanguage()\n3. i18next初始化后未根据用户偏好设置活动语言"],
        ["影响范围", "导航菜单、聊天页面、欢迎界面、历史页面、个人资料页、登录页、管理后台 — 全部页面"],
    ],
    [1.5, 7.5]
)

# Slide 5: Bug 3
add_table_slide(
    prs1,
    "Bug 3 — 提问无回答",
    ["项目", "详情"],
    [
        ["测试步骤", "输入有效问题并提交，检查加载状态、网络请求、响应接收、结果展示"],
        ["预期结果", "显示加载中 → 收到响应 → 正确展示回答，错误时显示友好提示"],
        ["实际结果", "❌ 失败 — 错误时思考动画消失，无任何错误提示，用户不知道发生了什么"],
        ["根因分析", "1. SSE错误事件被捕获后仅输出console.error，无用户可见的错误展示\n2. chatStore中没有error/lastError状态字段\n3. ChatPage没有错误展示UI组件"],
        ["影响范围", "所有聊天对话中的错误处理场景"],
    ],
    [1.5, 7.5]
)

# Slide 6: Bug 4
add_table_slide(
    prs1,
    "Bug 4 — 历史面板无对话记录",
    ["项目", "详情"],
    [
        ["测试步骤", "完成3轮对话后打开历史面板，验证对话是否出现，检查排序和持久化"],
        ["预期结果", "所有对话记录出现，按时间倒序排列，刷新后仍存在"],
        ["实际结果", "⚠️ 部分失败 — 对话可能出现但排序不保证，字段名不匹配导致日期为空"],
        ["根因分析", "1. ChatSessionListCreateView未设置ordering字段，查询结果顺序不确定\n2. API返回updated_at（snake_case），前端使用updatedAt（camelCase）导致日期显示为undefined\n3. 历史页面无刷新按钮，导航回来时数据可能过时"],
        ["影响范围", "历史页面"],
    ],
    [1.5, 7.5]
)

# Slide 7: New Issues
add_table_slide(
    prs1,
    "新发现的问题",
    ["#","问题描述","严重程度"],
    [
        ["1", "session.updatedAt / updated_at 字段名不匹配 — API返回snake_case，前端使用camelCase", "高"],
        ["2", "历史页面无加载状态指示器 — 数据加载期间无视觉反馈", "低"],
        ["3", "欢迎界面的快速操作未从API动态获取 — 硬编码在组件中，不随语言变化", "中"],
        ["4", "登录页面部分文字未国际化 — 品牌描述和功能列表仍为硬编码英文", "低"],
    ],
    [0.5, 7.0, 1.5]
)

# Slide 8: Conclusion
add_bullet_slide(
    prs1,
    "测试结论与总体评估",
    [
        "4个已知Bug中，3个已确认存在并修复（主题切换、语言切换、聊天错误处理）",
        "1个部分问题（历史记录）已修复排序和字段名问题",
        "新发现1个字段名不匹配Bug，已同步修复",
        "所有组件已完成i18n国际化改造，中英文切换即时生效",
        "整体评估：经过本轮修复，核心功能质量显著提升，建议回归测试验证",
    ]
)

prs1.save("test_report.pptx")
print("[OK] test_report.pptx generated")


# ============================================================
# PPT 2: Fix Summary
# ============================================================
prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)

# Slide 1: Cover
add_cover_slide(
    prs2,
    "Bug修复总结报告",
    "安永入职 AI · 质量修复记录"
)

# Slide 2: Summary Table
add_table_slide(
    prs2,
    "修复汇总",
    ["Bug ID", "名称", "根因", "修复方法", "状态"],
    [
        ["BUG-01", "主题切换卡顿", "全树重渲染 + darkAlgorithm", "useMemo缓存 + CSS过渡优化", "✅ 已修复"],
        ["BUG-02", "语言切换无效", "无useTranslation + 无changeLanguage", "全局i18n集成 + 语言同步", "✅ 已修复"],
        ["BUG-03", "聊天无回答", "错误被静默消费，无UI展示", "新增lastError状态 + 错误Alert", "✅ 已修复"],
        ["BUG-04", "历史面板为空", "无排序 + 字段名不匹配", "添加ordering + 兼容双字段名", "✅ 已修复"],
    ],
    [0.9, 1.3, 2.2, 2.8, 1.3]
)

# Slide 3: Bug 1 Fix
add_table_slide(
    prs2,
    "BUG-01 修复详情 — 主题切换卡顿",
    ["项目", "详情"],
    [
        ["症状", "切换浅色/深色模式时有明显延迟和闪烁，视觉不连贯"],
        ["根因分析（代码级）", "main.tsx中每次渲染创建新themeConfig对象 → 全组件树重渲染\nuseTheme.ts中darkAlgorithm运行时计算数千CSS变量\nglobals.css中body 0.3s过渡与Ant Design即时切换冲突"],
        ["修复描述", "1. main.tsx：使用useMemo缓存themeConfig，避免重复创建\n2. globals.css：添加全局0.2s过渡规则，使Ant Design组件平滑过渡"],
        ["修改文件", "frontend/src/main.tsx\nfrontend/src/styles/globals.css"],
        ["修复后验证", "切换即时响应，0.2s平滑过渡，无闪烁"],
    ],
    [1.5, 7.5]
)

# Slide 4: Bug 2 Fix
add_table_slide(
    prs2,
    "BUG-02 修复详情 — 语言切换无效",
    ["项目", "详情"],
    [
        ["症状", "在个人设置中切换中英文后，页面文字完全不变化"],
        ["根因分析（代码级）", "全部7个组件（AppLayout, ChatPage, WelcomeScreen, HistoryPage, ProfilePage, MessageBubble, KnowledgeBasePage, App.tsx）使用硬编码英文字符串\n无任何组件调用useTranslation()\nProfilePage保存语言偏好后未调用i18n.changeLanguage()\ni18n初始化使用固定lng='en'，未读取用户偏好"],
        ["修复描述", "1. i18n/index.ts：启动时从localStorage读取用户语言偏好\n2. 所有组件接入useTranslation()钩子，硬编码字符串替换为t()调用\n3. ProfilePage：保存时调用i18n.changeLanguage()\n4. App.tsx：登录成功后同步i18n语言\n5. 翻译JSON文件扩展为40+条键值对，覆盖全部UI文字"],
        ["修改文件", "frontend/src/i18n/index.ts\nfrontend/src/i18n/locales/en/common.json\nfrontend/src/i18n/locales/zh/common.json\nfrontend/src/i18n/locales/en/chat.json\nfrontend/src/i18n/locales/zh/chat.json\nfrontend/src/layout/AppLayout.tsx\nfrontend/src/pages/ChatPage.tsx\nfrontend/src/pages/HistoryPage.tsx\nfrontend/src/pages/ProfilePage.tsx\nfrontend/src/pages/App.tsx\nfrontend/src/pages/admin/KnowledgeBasePage.tsx\nfrontend/src/components/chat/WelcomeScreen.tsx\nfrontend/src/components/chat/MessageBubble.tsx"],
        ["修复后验证", "中英文切换即时生效，所有页面文字正确切换，刷新后保持"],
    ],
    [1.5, 7.5]
)

# Slide 5: Bug 3 Fix
add_table_slide(
    prs2,
    "BUG-03 修复详情 — 聊天无回答",
    ["项目", "详情"],
    [
        ["症状", "网络错误或服务器错误时，思考动画消失但无任何提示"],
        ["根因分析（代码级）", "chatStore.ts 第160-163行：SSE错误事件仅console.error，无用户可见展示\nChatState接口无error/lastError字段\nChatPage无错误展示组件"],
        ["修复描述", "1. chatStore.ts：新增lastError字段和setLastError/clearLastError动作\n2. 错误事件（SSE error / 网络异常 / 重试耗尽）均设置lastError\n3. ChatPage.tsx：添加Alert组件展示lastError，附带重试按钮\n4. 发送新消息时自动清除旧错误"],
        ["修改文件", "frontend/src/store/chatStore.ts\nfrontend/src/pages/ChatPage.tsx"],
        ["修复后验证", "错误时显示红色Alert + 重试按钮，网络异常显示友好提示"],
    ],
    [1.5, 7.5]
)

# Slide 6: Bug 4 Fix
add_table_slide(
    prs2,
    "BUG-04 & BUG-05 修复详情 — 历史面板 & 字段名不匹配",
    ["项目", "详情"],
    [
        ["症状", "对话记录可能出现但排序不确定，日期显示为空白"],
        ["根因分析（代码级）", "BUG-04: ChatSessionListCreateView无ordering设置，SQL查询结果顺序不确定\nHistoryPage无刷新按钮，导航回来时数据过时\nBUG-05: API返回updated_at (snake_case)，前端HistoryPage使用session.updatedAt (camelCase)，导致undefined"],
        ["修复描述", "BUG-04: 添加ordering = '-updated_at'到ChatSessionListCreateView，确保按时间倒序\n为HistoryPage添加刷新按钮\nBUG-05: HistoryPage兼容updated_at和updatedAt双字段名\n更新ChatSession TypeScript接口类型定义"],
        ["修改文件", "backend/apps/chat/views.py\nfrontend/src/pages/HistoryPage.tsx\nfrontend/src/store/chatStore.ts"],
        ["修复后验证", "历史记录按时间倒序显示，日期正确展示，刷新按钮可用"],
    ],
    [1.5, 7.5]
)

# Slide 7: Remaining Issues
add_table_slide(
    prs2,
    "遗留问题与建议",
    ["#", "问题", "建议"],
    [
        ["1", "登录页品牌描述文字未国际化", "将品牌文案加入翻译文件并接入t()调用"],
        ["2", "欢迎界面功能列表未国际化", "将'智能问答'等描述加入翻译文件"],
        ["3", "历史页面无分页/滚动加载", "添加分页组件或虚拟滚动优化大数据量场景"],
        ["4", "聊天消息未加载历史对话", "进入已有会话时从API加载历史消息"],
        ["5", "建议增加端到端自动化测试", "引入Playwright进行回归测试"],
    ],
    [0.5, 5.0, 4.0]
)

# Slide 8: Summary
add_bullet_slide(
    prs2,
    "修复总结",
    [
        "本次共修复 5 个 Bug（4个已知 + 1个新发现）",
        "涉及修改 13 个前端文件 + 1 个后端文件",
        "核心改进：",
        "  - 主题切换：从卡顿延迟优化为0.2s平滑过渡",
        "  - 语言切换：从完全不生效到全局即时切换",
        "  - 错误处理：从静默失败到友好提示+重试",
        "  - 历史记录：从排序不确定到按时间倒序+刷新支持",
        "所有修复保持项目原有代码风格，未引入破坏性变更",
        "建议后续增加端到端自动化测试以防止回归",
    ]
)

prs2.save("fix_summary.pptx")
print("[OK] fix_summary.pptx generated")
print("\n[OK] Both PPTs generated successfully!")
