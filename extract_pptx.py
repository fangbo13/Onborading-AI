# -*- coding: utf-8 -*-
from pptx import Presentation
with open('_pptx_output.txt', 'w', encoding='utf-8') as f:
    prs = Presentation('ux_audit_output/UX_Audit_Report_v3.pptx')
    for i, slide in enumerate(prs.slides):
        f.write(f'\n===== Slide {i+1} =====\n')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        f.write(text + '\n')
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    f.write(' | '.join(cells) + '\n')
print("Done. Output saved to _pptx_output.txt")
